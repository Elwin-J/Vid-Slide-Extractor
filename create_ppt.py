# kinda need this for importing pptx
# import collections
import collections.abc
import os
from PIL import Image

from pptx import Presentation


prs = Presentation()

# WIP
pth = './theframes'
blank_layout = prs.slide_layouts[6]
pics = os.listdir(pth)
pics.sort(key=lambda ele: int(ele.rstrip('.png').lstrip('frame')))
for pic in pics:
    impth = os.path.join(pth, pic)
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(impth, 0, 0, width=prs.slide_width, height=prs.slide_height)
    print(f'added slide {pic}')
prs.save("MyPresentation.pptx")